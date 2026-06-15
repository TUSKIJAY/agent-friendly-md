"""Format-native extraction routing probes and Phase 1 security audit.

These probes do not replace the vendored backend. They add evidence about what
the backend output can be trusted to represent, and they enrich normalized
elements only when the source file itself provides verifiable native structure.
"""
from __future__ import annotations

import json
import re
import zipfile
from collections import Counter
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

from . import extraction_elements as EE
from . import issues as ISS
from . import jobstate as js
from . import paths as P

FOUND_BY = "extraction_security_audit.py"
AUDIT_SCHEMA_VERSION = "agent-extraction-security-audit/0.1"
EVIDENCE_SCHEMA_VERSION = "agent-source-evidence-manifest/0.1"
PROMPT_PATTERNS = (
    re.compile(r"ignore\s+(all\s+)?previous\s+instructions", re.IGNORECASE),
    re.compile(r"disregard\s+(all\s+)?previous\s+instructions", re.IGNORECASE),
    re.compile(r"system\s+prompt", re.IGNORECASE),
    re.compile(r"prompt\s+injection", re.IGNORECASE),
    re.compile(r"忽略(以上|之前|前面).{0,12}指令"),
)
W_NS = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"


def analyze_and_write(
    job_root: str | Path,
    state: dict,
    meta: dict,
    elements: list[dict],
) -> tuple[list[dict], dict]:
    """Augment elements/meta with routing evidence and write audit artifacts."""
    root = Path(job_root)
    source_rel = (state.get("source_files") or [""])[0]
    source = root / source_rel
    source_type = EE.normalize_source_type(meta.get("source_format") or source.suffix)
    audit = _base_audit(source_rel, source_type)
    evidence_manifest = _base_evidence_manifest()
    enriched = [_clone(e) for e in elements]

    _mark_generic_fallbacks(enriched)
    _scan_prompt_injection(enriched, audit)

    routing: dict[str, Any] = {
        "schema_version": "agent-extraction-routing/0.1",
        "source_type": source_type,
        "route": "native_probe",
        "fallback_counts": _fallback_counts(enriched),
    }

    if source.is_file():
        if source_type == "pdf":
            routing["pdf"] = _augment_pdf(root, source, enriched, audit, evidence_manifest)
        elif source_type == "docx":
            routing["docx"] = _augment_docx(source, enriched, audit)
        elif source_type == "pptx":
            routing["pptx"] = _augment_pptx(source, enriched, audit)
        elif source_type == "xlsx":
            routing["xlsx"] = _augment_xlsx(source, enriched, audit)
        elif source_type in ("md", "txt"):
            routing[source_type] = _augment_text_source(source, enriched, audit)
    else:
        audit["findings"].append(_finding(
            "source_missing",
            "major",
            "security",
            f"Declared source file is missing: {source_rel}",
            "Restore the source file before trusting extraction evidence.",
            source_file=source_rel,
        ))

    routing["fallback_counts"] = _fallback_counts(enriched)
    meta = dict(meta)
    meta["routing"] = routing
    meta["evidence_level_stats"] = dict(sorted(Counter(e.get("evidence_level") for e in enriched).items()))
    meta["coverage"] = {**(meta.get("coverage") or {}), **_coverage_from_elements(enriched)}
    audit["routing_summary"] = routing
    audit["finding_count"] = len(audit["findings"])

    _write_audit(root, audit)
    _upsert_security_issues(root, audit)
    if evidence_manifest["evidence"]:
        _write_json(root / P.REVIEW_EVIDENCE_MANIFEST, evidence_manifest)
    return enriched, meta


def _base_audit(source_rel: str, source_type: str) -> dict:
    return {
        "schema_version": AUDIT_SCHEMA_VERSION,
        "created_at": js.now_iso(),
        "phase": "1_extraction",
        "source_file": source_rel,
        "source_type": source_type,
        "findings": [],
        "routing_summary": {},
        "finding_count": 0,
    }


def _base_evidence_manifest() -> dict:
    return {
        "schema_version": EVIDENCE_SCHEMA_VERSION,
        "created_at": js.now_iso(),
        "phase": "1_extraction",
        "evidence": [],
    }


def _clone(value: Any) -> Any:
    return json.loads(json.dumps(value, ensure_ascii=False))


def _write_json(path: Path, data: dict) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def _write_audit(job_root: Path, audit: dict) -> None:
    _write_json(job_root / P.REVIEW_EXTRACTION_SECURITY_AUDIT_JSON, audit)
    lines = [
        "# Extraction Security Audit",
        "",
        f"- schema_version: {audit['schema_version']}",
        f"- source_file: {audit['source_file']}",
        f"- source_type: {audit['source_type']}",
        f"- finding_count: {audit['finding_count']}",
        "",
        "## Findings",
        "",
    ]
    if not audit["findings"]:
        lines.append("_No extraction security findings._")
    for finding in audit["findings"]:
        lines.extend([
            f"### {finding['id']}",
            "",
            f"- severity: {finding['severity']}",
            f"- category: {finding['category']}",
            f"- issue: {finding['issue']}",
            f"- required_action: {finding['required_action']}",
            f"- source_file: {finding.get('source_file', '')}",
            f"- source_anchor: {finding.get('source_anchor', '')}",
            "",
        ])
    (job_root / P.REVIEW_EXTRACTION_SECURITY_AUDIT_MD).write_text(
        "\n".join(lines).rstrip() + "\n", encoding="utf-8"
    )


def _upsert_security_issues(job_root: Path, audit: dict) -> None:
    generated = [
        ISS.new_issue(
            issue_id=finding["id"],
            severity=finding["severity"],
            category="security",
            source_block_id="",
            source_anchor=str(finding.get("source_anchor") or ""),
            issue=finding["issue"],
            required_action=finding["required_action"],
            found_by=FOUND_BY,
        )
        for finding in audit["findings"]
        if finding.get("category") == "security" and finding.get("severity") in ("blocker", "major", "minor")
    ]
    issues, errors = ISS.upsert_generated(job_root / P.REVIEW_ISSUES, FOUND_BY, generated)
    if not errors:
        (job_root / P.REVIEW_UNRESOLVED).write_text(
            ISS.render_unresolved(ISS.unresolved_entries(issues)), encoding="utf-8"
        )


def _finding(
    finding_id: str,
    severity: str,
    category: str,
    issue: str,
    required_action: str,
    *,
    source_file: str = "",
    source_anchor: str = "",
    details: dict | None = None,
) -> dict:
    return {
        "id": f"security_{finding_id}",
        "severity": severity,
        "category": category,
        "issue": issue,
        "required_action": required_action,
        "source_file": source_file,
        "source_anchor": source_anchor,
        "details": details or {},
    }


def _metadata(element: dict) -> dict:
    meta = element.setdefault("native_metadata", {})
    if not isinstance(meta, dict):
        meta = {}
        element["native_metadata"] = meta
    return meta


def _remove_unavailable(meta: dict, *fields: str) -> None:
    unavailable = meta.get("unavailable")
    if not isinstance(unavailable, list):
        return
    remove = set(fields)
    meta["unavailable"] = [item for item in unavailable if item not in remove]


def _mark_generic_fallbacks(elements: list[dict]) -> None:
    for element in elements:
        meta = _metadata(element)
        flags = meta.setdefault("fallback", {})
        etype = element.get("element_type")
        if etype in ("image", "chart"):
            flags["needs_vlm"] = True
            flags["needs_agent_review"] = True
            element["needs_review"] = True
        elif element.get("needs_review"):
            flags["needs_agent_review"] = True


def _fallback_counts(elements: list[dict]) -> dict:
    counts = Counter()
    for element in elements:
        flags = (_metadata(element).get("fallback") or {})
        for key in ("needs_ocr", "needs_vlm", "needs_agent_review"):
            if flags.get(key):
                counts[key] += 1
    return dict(sorted(counts.items()))


def _coverage_from_elements(elements: list[dict]) -> dict:
    anchored = sum(1 for e in elements if e.get("source_anchor"))
    precise = 0
    for e in elements:
        anchor = e.get("source_anchor") or {}
        if any(anchor.get(k) for k in EE.HIGH_PRECISION_ANCHOR_FIELDS):
            precise += 1
    return {
        "elements_with_source_anchor": anchored,
        "elements_without_source_anchor": len(elements) - anchored,
        "elements_with_high_precision_anchor": precise,
        "elements_without_high_precision_anchor": len(elements) - precise,
    }


def _element_text(element: dict) -> str:
    content = element.get("content")
    if isinstance(content, str):
        return content
    if isinstance(content, dict):
        if "text" in content:
            return str(content.get("text") or "")
        if "items" in content and isinstance(content.get("items"), list):
            return "\n".join(str(x) for x in content["items"])
        if "table" in content:
            table = content.get("table") or {}
            parts = list(table.get("columns") or [])
            for row in table.get("rows") or []:
                for cell in row:
                    if isinstance(cell, dict):
                        parts.append(str(cell.get("text") or ""))
            return "\n".join(parts)
        if "figure" in content:
            fig = content.get("figure") or {}
            return "\n".join(str(fig.get(k) or "") for k in ("caption", "description", "asset_file"))
        if "latex" in content:
            return str(content.get("latex") or "")
        if "language" in content or "text" in content:
            return str(content.get("text") or "")
        return json.dumps(content, ensure_ascii=False)
    return ""


def _scan_prompt_injection(elements: list[dict], audit: dict) -> None:
    seen: set[str] = set()
    for element in elements:
        text = _element_text(element)
        if not text:
            continue
        for pattern in PROMPT_PATTERNS:
            if not pattern.search(text):
                continue
            anchor = element.get("source_anchor") or {}
            key = f"{element.get('element_id')}:{pattern.pattern}"
            if key in seen:
                continue
            seen.add(key)
            audit["findings"].append(_finding(
                f"prompt_injection_{element.get('element_id')}",
                "major",
                "security",
                "Extracted content contains text that resembles prompt-injection instructions.",
                "Review the source context and keep the text as quoted source content, not as operator instructions.",
                source_file=anchor.get("source_file", audit.get("source_file", "")),
                source_anchor=_anchor_label(anchor),
                details={"element_id": element.get("element_id"), "pattern": pattern.pattern},
            ))


def _anchor_label(anchor: dict) -> str:
    kind = anchor.get("kind", "")
    if kind == "pdf_page":
        return f"p.{anchor.get('page')}"
    if kind == "slide":
        return f"slide {anchor.get('slide')}"
    if kind == "sheet_range":
        return f"{anchor.get('sheet', '')}!{anchor.get('range', '')}".strip("!")
    if kind in ("docx_anchor", "text_anchor"):
        path = " / ".join(str(x) for x in anchor.get("heading_path") or [])
        idx = anchor.get("paragraph_index", anchor.get("block_index", ""))
        return f"{path} / block {idx}".strip(" /")
    return ""


def _augment_pdf(job_root: Path, source: Path, elements: list[dict], audit: dict, evidence_manifest: dict) -> dict:
    try:
        import fitz  # noqa: PLC0415
    except ImportError as e:
        audit["findings"].append(_finding(
            "pdf_probe_unavailable",
            "minor",
            "security",
            f"PyMuPDF unavailable, PDF native routing probe could not run: {e}",
            "Install PyMuPDF before relying on PDF text-layer coverage.",
            source_file=audit["source_file"],
        ))
        return {"status": "unavailable", "reason": str(e)}

    pages: list[dict] = []
    structured_blocks: list[dict] = []
    text_blocks_by_page: dict[int, list[dict]] = {}
    try:
        with fitz.open(source) as doc:
            for page_index, page in enumerate(doc, start=1):
                page_metrics, text_blocks = _pdf_page_metrics(page, page_index, audit["source_file"])
                for kind, item in page_metrics.pop("_findings", []):
                    label = "off-page" if kind == "off_page_text" else "very light"
                    audit["findings"].append(_finding(
                        f"pdf_{kind}_p{page_index}_{item['backend_element_id']}",
                        "major" if kind == "off_page_text" else "minor",
                        "security",
                        f"PDF text layer contains {label} text that may not be visible to readers.",
                        "Review the text-layer artifact and decide whether it is content, metadata, or injection risk.",
                        source_file=audit["source_file"],
                        source_anchor=f"p.{page_index}",
                        details=item,
                    ))
                page_metrics.pop("_source_file", None)
                pages.append(page_metrics)
                text_blocks_by_page[page_index] = text_blocks
                structured_blocks.extend(text_blocks)
                _write_pdf_overlay(job_root, page, page_metrics, text_blocks, evidence_manifest)
    except (OSError, RuntimeError, ValueError) as e:
        audit["findings"].append(_finding(
            "pdf_probe_failed",
            "major",
            "security",
            f"PDF native routing probe failed: {e}",
            "Repair or replace the PDF before trusting extraction evidence.",
            source_file=audit["source_file"],
        ))
        return {"status": "failed", "reason": str(e)}

    by_page = {p["page"]: p for p in pages}
    for element in elements:
        anchor = element.get("source_anchor") or {}
        if anchor.get("kind") != "pdf_page":
            continue
        page_no = anchor.get("page")
        page_metrics = by_page.get(page_no)
        if not page_metrics:
            continue
        meta = _metadata(element)
        routing = meta.setdefault("routing", {})
        route = page_metrics["route"]
        routing["pdf_page_route"] = route
        routing["pdf_text_chars"] = page_metrics["text_chars"]
        if route != "native_text":
            fallback = meta.setdefault("fallback", {})
            fallback["needs_ocr"] = True
            fallback["needs_agent_review"] = True
            element["needs_review"] = True
        if element.get("element_type") in ("heading", "paragraph", "list"):
            matched = _match_pdf_text_block(_element_text(element), text_blocks_by_page.get(page_no) or [])
            if matched:
                anchor["bbox"] = matched["bbox"]
                anchor["backend_element_id"] = matched["backend_element_id"]
                anchor["layout_zone"] = matched["layout_zone"]
                meta.setdefault("native_probe", {})["pdf_text_layer_match"] = "unique_text_block"
                _remove_unavailable(
                    meta,
                    "source_anchor.bbox",
                    "source_anchor.backend_element_id",
                    "source_anchor.layout_zone",
                )
    text_blocks_file = _write_pdf_text_blocks(job_root, structured_blocks)
    return {
        "status": "ok",
        "page_count": len(pages),
        "pages": pages,
        "native_text_pages": sum(1 for p in pages if p["route"] == "native_text"),
        "fallback_pages": [p["page"] for p in pages if p["route"] != "native_text"],
        "text_blocks_file": text_blocks_file,
    }


def _pdf_page_metrics(page: Any, page_index: int, source_file: str) -> tuple[dict, list[dict]]:
    rect = page.rect
    data = page.get_text("dict")
    text_blocks: list[dict] = []
    off_page: list[dict] = []
    invisible_like: list[dict] = []
    text_chars = 0
    for idx, block in enumerate(data.get("blocks") or [], start=1):
        if block.get("type") != 0:
            continue
        text = _pdf_block_text(block)
        if not text.strip():
            continue
        bbox = [float(x) for x in block.get("bbox", (0, 0, 0, 0))]
        clipped = _clip_bbox(bbox, rect)
        block_area = max(0.0, (bbox[2] - bbox[0]) * (bbox[3] - bbox[1]))
        clipped_area = max(0.0, (clipped[2] - clipped[0]) * (clipped[3] - clipped[1]))
        backend_id = f"p{page_index}-textblock-{idx}"
        record = {
            "backend_element_id": backend_id,
            "text": text,
            "bbox": [round(x, 2) for x in bbox],
            "layout_zone": _layout_zone(bbox, rect),
        }
        text_blocks.append(record)
        text_chars += len(text.strip())
        if block_area and clipped_area / block_area < 0.5:
            off_page.append({"backend_element_id": backend_id, "bbox": record["bbox"], "text_preview": text[:80]})
        if _block_is_very_light(block):
            invisible_like.append({"backend_element_id": backend_id, "bbox": record["bbox"], "text_preview": text[:80]})
    route = "native_text" if text_chars >= 8 else "ocr_or_vlm_fallback"
    metrics = {
        "page": page_index,
        "text_chars": text_chars,
        "text_block_count": len(text_blocks),
        "image_count": len(page.get_images(full=True)),
        "route": route,
        "off_page_text_count": len(off_page),
        "very_light_text_count": len(invisible_like),
    }
    findings = []
    for item in off_page:
        findings.append(("off_page_text", item))
    for item in invisible_like:
        findings.append(("very_light_text", item))
    metrics["_findings"] = findings
    metrics["_source_file"] = source_file
    return metrics, text_blocks


def _pdf_block_text(block: dict) -> str:
    parts: list[str] = []
    for line in block.get("lines") or []:
        for span in line.get("spans") or []:
            parts.append(str(span.get("text") or ""))
        parts.append("\n")
    return "".join(parts).strip()


def _clip_bbox(bbox: list[float], rect: Any) -> list[float]:
    return [
        max(float(rect.x0), min(float(rect.x1), bbox[0])),
        max(float(rect.y0), min(float(rect.y1), bbox[1])),
        max(float(rect.x0), min(float(rect.x1), bbox[2])),
        max(float(rect.y0), min(float(rect.y1), bbox[3])),
    ]


def _layout_zone(bbox: list[float], rect: Any) -> str:
    center_y = (bbox[1] + bbox[3]) / 2
    height = float(rect.height or 1)
    if center_y < height * 0.2:
        return "header"
    if center_y > height * 0.85:
        return "footer"
    return "body"


def _block_is_very_light(block: dict) -> bool:
    colors = []
    for line in block.get("lines") or []:
        for span in line.get("spans") or []:
            color = span.get("color")
            if isinstance(color, int):
                colors.append(color)
    return bool(colors) and all(color >= 0xF0F0F0 for color in colors)


def _match_pdf_text_block(text: str, blocks: list[dict]) -> dict | None:
    needle = _norm_text(text)
    if len(needle) < 3:
        return None
    matches = [block for block in blocks if needle in _norm_text(block.get("text", ""))]
    if len(matches) == 1:
        return matches[0]
    return None


def _norm_text(text: str) -> str:
    return re.sub(r"\s+", " ", str(text)).strip().lower()


def _write_pdf_overlay(
    job_root: Path,
    page: Any,
    page_metrics: dict,
    text_blocks: list[dict],
    evidence_manifest: dict,
) -> None:
    if not text_blocks:
        return
    try:
        from PIL import Image, ImageDraw  # noqa: PLC0415
    except ImportError:
        return
    out = (
        job_root / P.REVIEW_VISUAL / "extraction_overlays"
        / "pdf_page_bbox_overlay" / f"page_{page_metrics['page']:03d}_overlay.png"
    )
    out.parent.mkdir(parents=True, exist_ok=True)
    # Keep scale 1 so bbox coordinates align with page coordinates.
    pix = page.get_pixmap(alpha=False)
    img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
    draw = ImageDraw.Draw(img)
    for block in text_blocks:
        bbox = block["bbox"]
        draw.rectangle(bbox, outline=(220, 40, 40), width=2)
        draw.text((bbox[0], max(0, bbox[1] - 10)), block["backend_element_id"], fill=(40, 40, 40))
    img.save(out)
    evidence_manifest["evidence"].append({
        "type": "extraction_overlay",
        "source_anchor": {"kind": "pdf_page", "page": page_metrics["page"]},
        "file": out.relative_to(job_root).as_posix(),
        "renderer_ok": True,
        "producer": "pdf_text_layer_probe",
    })


def _write_pdf_text_blocks(job_root: Path, blocks: list[dict]) -> str | None:
    if not blocks:
        return None
    path = job_root / P.EXTRACT_PDF_TEXT_BLOCKS
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8") as fh:
        for block in blocks:
            fh.write(json.dumps(block, ensure_ascii=False) + "\n")
    return P.EXTRACT_PDF_TEXT_BLOCKS


def _augment_docx(source: Path, elements: list[dict], audit: dict) -> dict:
    try:
        from docx import Document  # noqa: PLC0415
    except ImportError as e:
        return {"status": "unavailable", "reason": str(e)}
    doc = Document(str(source))
    profile = _docx_ooxml_profile(source)
    if profile["hidden_text_count"] or profile["tracked_deletion_count"]:
        audit["findings"].append(_finding(
            "docx_hidden_or_deleted_text",
            "minor",
            "security",
            "DOCX contains hidden text or tracked deletion markup.",
            "Review OOXML hidden/deleted text before treating extracted body text as complete.",
            source_file=audit["source_file"],
            source_anchor="docx body",
            details={
                "hidden_text_count": profile["hidden_text_count"],
                "tracked_deletion_count": profile["tracked_deletion_count"],
            },
        ))
    table_index = 0
    image_elements = [
        e for e in elements
        if e.get("element_type") == "image" and (e.get("source_anchor") or {}).get("kind") == "docx_anchor"
    ]
    image_rel_ids = list(profile["image_rel_ids"])
    image_rel_cursor = 0
    one_to_one_images = bool(image_rel_ids) and len(image_rel_ids) == len(image_elements)
    comment_elements = [
        e for e in elements
        if e.get("element_type") == "comment" and (e.get("source_anchor") or {}).get("kind") == "docx_anchor"
    ]
    for element in elements:
        anchor = element.get("source_anchor") or {}
        if anchor.get("kind") != "docx_anchor":
            continue
        meta = _metadata(element)
        meta.setdefault("routing", {})["docx_native"] = "ooxml_flow"
        if element.get("element_type") == "table":
            table_index += 1
            anchor.setdefault("table_index", table_index)
            _remove_unavailable(meta, "source_anchor.table_index")
        elif element.get("element_type") == "image" and one_to_one_images:
            anchor["image_rel_id"] = image_rel_ids[image_rel_cursor]
            image_rel_cursor += 1
            meta.setdefault("native_probe", {})["docx_image_match"] = "relationship_order_one_to_one"
            _remove_unavailable(meta, "source_anchor.image_rel_id")
        elif element.get("element_type") == "comment":
            matched = _match_docx_comment(_element_text(element), profile["comments"])
            if matched is None and len(comment_elements) == 1 and len(profile["comments"]) == 1:
                matched = profile["comments"][0]
            if matched:
                anchor["comment_id"] = matched["comment_id"]
                meta.setdefault("native_probe", {})["docx_comment_match"] = matched["match"]
                _remove_unavailable(meta, "source_anchor.comment_id")
    return {
        "status": "ok",
        "paragraph_count": len(doc.paragraphs),
        "table_count": len(doc.tables),
        "inline_shape_count": len(doc.inline_shapes),
        "image_relationship_count": len(profile["image_rel_ids"]),
        "comment_count": len(profile["comments"]),
        "footnote_count": profile["footnote_count"],
        "endnote_count": profile["endnote_count"],
        "hidden_text_count": profile["hidden_text_count"],
        "tracked_deletion_count": profile["tracked_deletion_count"],
    }


def _docx_ooxml_profile(source: Path) -> dict:
    profile = {
        "image_rel_ids": [],
        "comments": [],
        "footnote_count": 0,
        "endnote_count": 0,
        "hidden_text_count": 0,
        "tracked_deletion_count": 0,
    }
    try:
        with zipfile.ZipFile(source) as zf:
            names = set(zf.namelist())
            if "word/_rels/document.xml.rels" in names:
                rel_root = ET.fromstring(zf.read("word/_rels/document.xml.rels"))
                for rel in rel_root:
                    rel_type = rel.attrib.get("Type", "")
                    target = rel.attrib.get("Target", "")
                    rel_id = rel.attrib.get("Id", "")
                    if rel_id and (rel_type.endswith("/image") or target.startswith("media/")):
                        profile["image_rel_ids"].append(rel_id)
            if "word/document.xml" in names:
                doc_root = ET.fromstring(zf.read("word/document.xml"))
                embedded = [
                    blip.attrib.get(f"{R_NS}embed")
                    for blip in doc_root.iter()
                    if str(blip.tag).endswith("}blip") and blip.attrib.get(f"{R_NS}embed")
                ]
                for rel_id in embedded:
                    if rel_id not in profile["image_rel_ids"]:
                        profile["image_rel_ids"].append(rel_id)
                profile["hidden_text_count"] = sum(1 for el in doc_root.iter() if el.tag == f"{W_NS}vanish")
                profile["tracked_deletion_count"] = (
                    sum(1 for el in doc_root.iter() if el.tag == f"{W_NS}del")
                    + sum(1 for el in doc_root.iter() if el.tag == f"{W_NS}delText")
                )
            if "word/comments.xml" in names:
                comments_root = ET.fromstring(zf.read("word/comments.xml"))
                for comment in comments_root.iter(f"{W_NS}comment"):
                    cid = comment.attrib.get(f"{W_NS}id") or comment.attrib.get("id")
                    text = "".join(t.text or "" for t in comment.iter(f"{W_NS}t")).strip()
                    if cid is not None:
                        profile["comments"].append({
                            "comment_id": str(cid),
                            "text": text,
                            "match": "comments_xml_text",
                        })
            for rel, key in (("word/footnotes.xml", "footnote_count"), ("word/endnotes.xml", "endnote_count")):
                if rel in names:
                    root = ET.fromstring(zf.read(rel))
                    tag = f"{W_NS}{'footnote' if 'foot' in rel else 'endnote'}"
                    profile[key] = sum(1 for el in root.iter(tag))
    except (OSError, zipfile.BadZipFile, ET.ParseError):
        return profile
    return profile


def _match_docx_comment(text: str, comments: list[dict]) -> dict | None:
    needle = _norm_text(text)
    if not needle:
        return None
    matches = [c for c in comments if needle in _norm_text(c.get("text", ""))]
    return matches[0] if len(matches) == 1 else None


def _augment_pptx(source: Path, elements: list[dict], audit: dict) -> dict:
    try:
        from pptx import Presentation  # noqa: PLC0415
    except ImportError as e:
        return {"status": "unavailable", "reason": str(e)}
    prs = Presentation(str(source))
    slide_profiles: dict[int, dict] = {}
    text_shapes: dict[int, list[dict]] = {}
    table_shapes: dict[int, list[dict]] = {}
    chart_shapes: dict[int, list[dict]] = {}
    picture_shapes: dict[int, list[dict]] = {}
    note_shapes: dict[int, list[dict]] = {}
    for slide_no, slide in enumerate(prs.slides, start=1):
        shapes = []
        for z_order, shape in enumerate(slide.shapes, start=1):
            rec = _pptx_shape_record(shape, z_order, prs.slide_width, prs.slide_height)
            shapes.append(rec)
            if rec["off_slide"]:
                audit["findings"].append(_finding(
                    f"pptx_off_slide_shape_s{slide_no}_{rec['shape_id']}",
                    "minor",
                    "security",
                    "PPTX shape is partially or fully outside the slide canvas.",
                    "Review off-slide content before treating extracted slide text as complete.",
                    source_file=audit["source_file"],
                    source_anchor=f"slide {slide_no}",
                    details={k: rec[k] for k in ("shape_id", "shape_name", "bbox")},
                ))
        notes = _pptx_notes_text(slide)
        if notes:
            note_shapes[slide_no] = [{
                "shape_id": f"slide_{slide_no}_notes",
                "shape_name": "speaker_notes",
                "z_order": None,
                "text": notes,
                "bbox": None,
                "placeholder_type": "speaker_notes",
                "has_table": False,
                "has_chart": False,
                "has_picture": False,
                "table_range": None,
                "off_slide": False,
            }]
            audit["findings"].append(_finding(
                f"pptx_speaker_notes_s{slide_no}",
                "minor",
                "security",
                "PPTX slide contains speaker notes that may not appear in rendered slide screenshots.",
                "Review speaker notes and decide whether they are source content or presenter-only context.",
                source_file=audit["source_file"],
                source_anchor=f"slide {slide_no}",
                details={"text_preview": notes[:120]},
            ))
        else:
            note_shapes[slide_no] = []
        slide_profiles[slide_no] = {
            "shape_count": len(shapes),
            "table_shape_count": sum(1 for s in shapes if s["has_table"]),
            "chart_shape_count": sum(1 for s in shapes if s["has_chart"]),
            "picture_shape_count": sum(1 for s in shapes if s["has_picture"]),
            "off_slide_shape_count": sum(1 for s in shapes if s["off_slide"]),
            "speaker_notes_chars": len(notes),
        }
        text_shapes[slide_no] = [s for s in shapes if s["text"].strip()]
        table_shapes[slide_no] = [s for s in shapes if s["has_table"]]
        chart_shapes[slide_no] = [s for s in shapes if s["has_chart"]]
        picture_shapes[slide_no] = [s for s in shapes if s["has_picture"]]
    image_positions = _sequence_positions(elements, "image", "slide")
    chart_positions = _sequence_positions(elements, "chart", "slide")
    for element in elements:
        anchor = element.get("source_anchor") or {}
        if anchor.get("kind") != "slide":
            continue
        slide_no = anchor.get("slide")
        meta = _metadata(element)
        meta.setdefault("routing", {})["pptx_native"] = "slide_shape_tree"
        matched = None
        etype = element.get("element_type")
        if etype == "table":
            matched = _match_pptx_table(element, table_shapes.get(slide_no) or [])
        elif etype == "chart":
            matched = _match_single_by_sequence(element, chart_shapes.get(slide_no) or [], chart_positions)
        elif etype == "image":
            matched = _match_single_by_sequence(element, picture_shapes.get(slide_no) or [], image_positions)
        elif etype == "note":
            matched = _match_pptx_shape(_element_text(element), note_shapes.get(slide_no) or [])
        else:
            matched = _match_pptx_shape(_element_text(element), text_shapes.get(slide_no) or [])
        if matched:
            _apply_pptx_shape_anchor(anchor, meta, matched)
            if etype == "table" and matched.get("table_range"):
                anchor["cell_ref"] = matched["table_range"]
                _remove_unavailable(meta, "source_anchor.cell_ref")
        if etype in ("image", "chart"):
            meta.setdefault("fallback", {})["needs_vlm"] = True
            meta.setdefault("fallback", {})["needs_agent_review"] = True
            element["needs_review"] = True
    return {"status": "ok", "slide_count": len(prs.slides), "slides": slide_profiles}


def _pptx_shape_record(shape: Any, z_order: int, slide_width: int, slide_height: int) -> dict:
    text = str(getattr(shape, "text", "") or "")
    has_table = bool(getattr(shape, "has_table", False))
    table_cells = _pptx_table_cells(shape) if has_table else []
    bbox = _pptx_shape_bbox(shape)
    return {
        "shape_id": str(shape.shape_id),
        "shape_name": str(shape.name or ""),
        "z_order": z_order,
        "text": text,
        "bbox": bbox,
        "placeholder_type": _pptx_placeholder_type(shape),
        "has_table": has_table,
        "has_chart": bool(getattr(shape, "has_chart", False)),
        "has_picture": _pptx_shape_has_blip(shape),
        "table_text": "\n".join(cell["text"] for cell in table_cells),
        "table_range": _pptx_table_range(table_cells),
        "off_slide": _pptx_off_slide(bbox, slide_width, slide_height),
    }


def _pptx_shape_bbox(shape: Any) -> list[int] | None:
    values = []
    for attr in ("left", "top", "width", "height"):
        value = getattr(shape, attr, None)
        if value is None:
            return None
        values.append(int(value))
    return values


def _pptx_placeholder_type(shape: Any) -> str | None:
    try:
        if getattr(shape, "is_placeholder", False):
            return str(shape.placeholder_format.type)
    except (AttributeError, ValueError):
        return None
    return None


def _pptx_shape_has_blip(shape: Any) -> bool:
    try:
        return any(str(el.tag).endswith("}blip") for el in shape._element.iter())
    except AttributeError:
        return False


def _pptx_off_slide(bbox: list[int] | None, slide_width: int, slide_height: int) -> bool:
    if not bbox:
        return False
    left, top, width, height = bbox
    return left < 0 or top < 0 or left + width > int(slide_width) or top + height > int(slide_height)


def _pptx_table_cells(shape: Any) -> list[dict]:
    cells: list[dict] = []
    try:
        rows = shape.table.rows
        cols = shape.table.columns
    except AttributeError:
        return cells
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, cell in enumerate(row.cells, start=1):
            text = str(cell.text or "").strip()
            if text:
                cells.append({"row": r_idx, "col": c_idx, "text": text})
    if not cells:
        return []
    return cells + [{"row_count": len(rows), "col_count": len(cols), "text": ""}]


def _pptx_table_range(cells: list[dict]) -> str | None:
    dims = next((c for c in cells if "row_count" in c), None)
    if not dims:
        return None
    return f"R1C1:R{dims['row_count']}C{dims['col_count']}"


def _pptx_notes_text(slide: Any) -> str:
    try:
        if not getattr(slide, "has_notes_slide", False):
            return ""
        text_frame = slide.notes_slide.notes_text_frame
    except (AttributeError, ValueError):
        return ""
    if text_frame is None:
        return ""
    return "\n".join(p.text for p in text_frame.paragraphs if p.text).strip()


def _match_pptx_shape(text: str, shapes: list[dict]) -> dict | None:
    needle = _norm_text(text)
    if len(needle) < 2:
        return None
    matches = [shape for shape in shapes if needle in _norm_text(shape["text"])]
    return matches[0] if len(matches) == 1 else None


def _match_pptx_table(element: dict, shapes: list[dict]) -> dict | None:
    if not shapes:
        return None
    table_text = _norm_text(_table_text_from_element(element))
    if table_text:
        matches = [shape for shape in shapes if table_text in _norm_text(shape.get("table_text", ""))]
        if len(matches) == 1:
            return matches[0]
    return shapes[0] if len(shapes) == 1 else None


def _sequence_positions(elements: list[dict], element_type: str, anchor_kind: str) -> dict[str, int]:
    positions: dict[str, int] = {}
    index_by_scope: Counter = Counter()
    for element in elements:
        anchor = element.get("source_anchor") or {}
        if element.get("element_type") != element_type or anchor.get("kind") != anchor_kind:
            continue
        scope = _anchor_scope(anchor)
        key = f"{scope}:{element.get('element_id')}"
        index_by_scope[scope] += 1
        positions[key] = index_by_scope[scope]
    return positions


def _anchor_scope(anchor: dict) -> Any:
    if anchor.get("kind") == "slide":
        return anchor.get("slide")
    if anchor.get("kind") == "sheet_range":
        return anchor.get("sheet")
    return anchor.get("source_file")


def _match_single_by_sequence(element: dict, shapes: list[dict], positions: dict[str, int]) -> dict | None:
    if not shapes:
        return None
    anchor = element.get("source_anchor") or {}
    if len(shapes) == 1:
        return shapes[0]
    pos = positions.get(f"{anchor.get('slide')}:{element.get('element_id')}")
    if pos and pos <= len(shapes):
        return shapes[pos - 1]
    return None


def _apply_pptx_shape_anchor(anchor: dict, meta: dict, matched: dict) -> None:
    anchor["shape_id"] = matched["shape_id"]
    anchor["shape_name"] = matched["shape_name"]
    if matched.get("z_order") is not None:
        anchor["z_order"] = matched["z_order"]
    if matched.get("bbox") is not None:
        anchor["bbox"] = matched["bbox"]
    if matched.get("placeholder_type"):
        anchor["placeholder_type"] = matched["placeholder_type"]
    _remove_unavailable(
        meta,
        "source_anchor.shape_id",
        "source_anchor.shape_name",
        "source_anchor.z_order",
        "source_anchor.bbox",
        "source_anchor.placeholder_type",
    )


def _augment_xlsx(source: Path, elements: list[dict], audit: dict) -> dict:
    try:
        import openpyxl  # noqa: PLC0415
        from openpyxl.utils import get_column_letter, range_boundaries  # noqa: PLC0415
    except ImportError as e:
        return {"status": "unavailable", "reason": str(e)}
    wb = openpyxl.load_workbook(source, data_only=False, read_only=False)
    sheets = {}
    tables_by_sheet: dict[str, list[dict]] = {}
    formulas_by_sheet: dict[str, list[dict]] = {}
    charts_by_sheet: dict[str, list[dict]] = {}
    for ws in wb.worksheets:
        dimension = ws.calculate_dimension()
        hidden_rows = [idx for idx, dim in ws.row_dimensions.items() if dim.hidden]
        hidden_cols = [key for key, dim in ws.column_dimensions.items() if dim.hidden]
        tables = _xlsx_tables(ws, range_boundaries)
        formulas = _xlsx_formulas(ws)
        charts = _xlsx_charts(ws, get_column_letter)
        merged_ranges = [str(rng) for rng in ws.merged_cells.ranges]
        tables_by_sheet[ws.title] = tables
        formulas_by_sheet[ws.title] = formulas
        charts_by_sheet[ws.title] = charts
        sheets[ws.title] = {
            "range": dimension,
            "sheet_state": ws.sheet_state,
            "hidden_row_count": len(hidden_rows),
            "hidden_column_count": len(hidden_cols),
            "table_count": len(tables),
            "formula_cell_count": len(formulas),
            "chart_count": len(charts),
            "merged_range_count": len(merged_ranges),
            "image_count": len(getattr(ws, "_images", []) or []),
            "tables": [{"name": t["table_name"], "range": t["range"]} for t in tables],
            "formula_cells": [f["formula_cell"] for f in formulas],
            "charts": [{"chart_id": c["chart_id"], "range": c.get("range")} for c in charts],
            "merged_ranges": merged_ranges,
        }
        if ws.sheet_state != "visible" or hidden_rows or hidden_cols:
            audit["findings"].append(_finding(
                f"xlsx_hidden_structure_{_slug(ws.title)}",
                "minor",
                "security",
                "Workbook contains hidden sheet, row, or column structure.",
                "Review hidden workbook structure before treating the extracted table view as complete.",
                source_file=audit["source_file"],
                source_anchor=f"{ws.title}!{dimension}",
                details=sheets[ws.title],
            ))
        if merged_ranges:
            audit["findings"].append(_finding(
                f"xlsx_merged_ranges_{_slug(ws.title)}",
                "minor",
                "structure",
                "Worksheet contains merged cell ranges that can affect table interpretation.",
                "Review merged ranges during table QA.",
                source_file=audit["source_file"],
                source_anchor=f"{ws.title}!{dimension}",
                details={"merged_ranges": merged_ranges},
            ))
    chart_positions = _sequence_positions(elements, "chart", "sheet_range")
    for element in elements:
        anchor = element.get("source_anchor") or {}
        if anchor.get("kind") != "sheet_range":
            continue
        sheet = anchor.get("sheet")
        profile = sheets.get(sheet)
        if not profile:
            continue
        meta = _metadata(element)
        meta.setdefault("routing", {})["xlsx_native"] = "workbook_model"
        etype = element.get("element_type")
        if etype == "table":
            table = _match_xlsx_table(element, tables_by_sheet.get(sheet) or [])
            if table:
                anchor["range"] = table["range"]
                anchor["table_name"] = table["table_name"]
                _remove_unavailable(meta, "source_anchor.range", "source_anchor.table_name")
            elif not anchor.get("range"):
                anchor["range"] = profile["range"]
                _remove_unavailable(meta, "source_anchor.range")
        elif etype == "formula":
            formula = _match_xlsx_formula(element, formulas_by_sheet.get(sheet) or [])
            if formula:
                anchor["formula_cell"] = formula["formula_cell"]
                anchor["range"] = formula["formula_cell"]
                _remove_unavailable(meta, "source_anchor.formula_cell", "source_anchor.range")
        elif etype == "chart":
            chart = _match_xlsx_by_sequence(element, charts_by_sheet.get(sheet) or [], chart_positions)
            if chart:
                anchor["chart_id"] = chart["chart_id"]
                if chart.get("range"):
                    anchor["range"] = chart["range"]
                    _remove_unavailable(meta, "source_anchor.range")
                _remove_unavailable(meta, "source_anchor.chart_id")
            meta.setdefault("fallback", {})["needs_vlm"] = True
            meta.setdefault("fallback", {})["needs_agent_review"] = True
            element["needs_review"] = True
        elif etype == "image":
            image_count = int(profile.get("image_count") or 0)
            if image_count:
                meta.setdefault("native_probe", {})["xlsx_image_count"] = image_count
            meta.setdefault("fallback", {})["needs_vlm"] = True
            meta.setdefault("fallback", {})["needs_agent_review"] = True
            element["needs_review"] = True
        elif not anchor.get("range"):
            anchor["range"] = profile["range"]
            _remove_unavailable(meta, "source_anchor.range")
    wb.close()
    return {"status": "ok", "sheet_count": len(sheets), "sheets": sheets}


def _xlsx_tables(ws: Any, range_boundaries: Any) -> list[dict]:
    tables = []
    table_names = list(getattr(ws, "tables", {}) or {})
    for name in table_names:
        table_obj = ws.tables[name]
        table_name = str(getattr(table_obj, "displayName", None) or getattr(table_obj, "name", None) or name)
        ref = str(getattr(table_obj, "ref", None) or table_obj)
        text = []
        try:
            min_col, min_row, max_col, max_row = range_boundaries(ref)
            for row in ws.iter_rows(min_row=min_row, max_row=max_row, min_col=min_col, max_col=max_col):
                for cell in row:
                    if cell.value is not None:
                        text.append(str(cell.value))
        except (TypeError, ValueError):
            pass
        tables.append({"table_name": table_name, "range": ref, "text": "\n".join(text)})
    return tables


def _xlsx_formulas(ws: Any) -> list[dict]:
    formulas = []
    for row in ws.iter_rows():
        for cell in row:
            value = cell.value
            if cell.data_type == "f" or (isinstance(value, str) and value.startswith("=")):
                formulas.append({"formula_cell": cell.coordinate, "formula": str(value)})
    return formulas


def _xlsx_charts(ws: Any, get_column_letter: Any) -> list[dict]:
    charts = []
    for idx, chart in enumerate(getattr(ws, "_charts", []) or [], start=1):
        cell = _xlsx_anchor_cell(getattr(chart, "anchor", None), get_column_letter)
        charts.append({"chart_id": f"chart_{idx}", "range": cell, "title": str(getattr(chart, "title", "") or "")})
    return charts


def _xlsx_anchor_cell(anchor: Any, get_column_letter: Any) -> str | None:
    marker = getattr(anchor, "_from", None)
    if marker is None:
        return None
    try:
        return f"{get_column_letter(int(marker.col) + 1)}{int(marker.row) + 1}"
    except (TypeError, ValueError):
        return None


def _table_text_from_element(element: dict) -> str:
    content = element.get("content")
    if not isinstance(content, dict):
        return _element_text(element)
    table = content.get("table") or {}
    parts = [str(c) for c in table.get("columns") or []]
    for row in table.get("rows") or []:
        for cell in row:
            if isinstance(cell, dict):
                parts.append(str(cell.get("text") or ""))
            else:
                parts.append(str(cell))
    return "\n".join(p for p in parts if p)


def _match_xlsx_table(element: dict, tables: list[dict]) -> dict | None:
    if not tables:
        return None
    needle = _norm_text(_table_text_from_element(element))
    if needle:
        matches = [table for table in tables if needle in _norm_text(table.get("text", ""))]
        if len(matches) == 1:
            return matches[0]
    return tables[0] if len(tables) == 1 else None


def _match_xlsx_formula(element: dict, formulas: list[dict]) -> dict | None:
    if not formulas:
        return None
    needle = _norm_text(_element_text(element)).strip("$")
    if needle:
        matches = [
            formula for formula in formulas
            if needle in _norm_text(formula.get("formula", "")) or _norm_text(formula.get("formula", "")) in needle
        ]
        if len(matches) == 1:
            return matches[0]
    return formulas[0] if len(formulas) == 1 else None


def _match_xlsx_by_sequence(element: dict, items: list[dict], positions: dict[str, int]) -> dict | None:
    if not items:
        return None
    anchor = element.get("source_anchor") or {}
    if len(items) == 1:
        return items[0]
    pos = positions.get(f"{anchor.get('sheet')}:{element.get('element_id')}")
    if pos and pos <= len(items):
        return items[pos - 1]
    return None


def _augment_text_source(source: Path, elements: list[dict], audit: dict) -> dict:
    text = source.read_text(encoding="utf-8", errors="replace")
    lines = text.splitlines()
    cursor = 0
    matched = 0
    for element in elements:
        anchor = element.get("source_anchor") or {}
        if anchor.get("kind") != "text_anchor":
            continue
        needle = _first_search_line(element)
        if not needle:
            continue
        hit = _find_line(lines, needle, cursor)
        if hit is None:
            continue
        anchor["line_range"] = [hit + 1, hit + 1]
        cursor = hit + 1
        matched += 1
        _remove_unavailable(_metadata(element), "source_anchor.line_range")
    return {"status": "ok", "line_count": len(lines), "line_range_matched_elements": matched}


def _first_search_line(element: dict) -> str:
    text = _element_text(element)
    for line in text.splitlines():
        line = line.strip().strip("#").strip()
        if len(line) >= 2:
            return line[:120]
    return ""


def _find_line(lines: list[str], needle: str, start: int) -> int | None:
    needle_norm = _norm_text(needle)
    for idx in range(start, len(lines)):
        if needle_norm and needle_norm in _norm_text(lines[idx]):
            return idx
    for idx, line in enumerate(lines):
        if needle_norm and needle_norm in _norm_text(line):
            return idx
    return None


def _slug(text: str) -> str:
    return re.sub(r"[^a-z0-9_]+", "_", text.lower()).strip("_")[:40] or "item"
